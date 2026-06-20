"""
Universal Document Extractor Service
Mendukung: PDF, Word (.doc/.docx), Excel (.xls/.xlsx), PowerPoint (.ppt/.pptx),
           Text (.txt/.csv/.tsv), Gambar (.jpg/.jpeg/.png/.webp/.bmp/.tiff), Struk/Nota
"""
import os
import uuid
import asyncio
import aiofiles
import logging
from datetime import datetime, timezone
from pathlib import Path
from typing import Optional
import re

from app.core.config import settings

logger = logging.getLogger(__name__)


# ─────────────────────────────────────────────
# MIME type & ekstensi yang didukung
# ─────────────────────────────────────────────

SUPPORTED_EXTENSIONS: dict[str, str] = {
    # PDF
    ".pdf":  "application/pdf",
    # Word
    ".doc":  "application/msword",
    ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    # Excel
    ".xls":  "application/vnd.ms-excel",
    ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    # PowerPoint
    ".ppt":  "application/vnd.ms-powerpoint",
    ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    # Text
    ".txt":  "text/plain",
    ".csv":  "text/csv",
    ".tsv":  "text/tab-separated-values",
    # Gambar
    ".jpg":  "image/jpeg",
    ".jpeg": "image/jpeg",
    ".png":  "image/png",
    ".webp": "image/webp",
    ".bmp":  "image/bmp",
    ".tiff": "image/tiff",
    ".tif":  "image/tiff",
}

IMAGE_EXTENSIONS  = {".jpg", ".jpeg", ".png", ".webp", ".bmp", ".tiff", ".tif"}
TEXT_EXTENSIONS   = {".txt", ".csv", ".tsv"}
PDF_EXTENSIONS    = {".pdf"}
WORD_EXTENSIONS   = {".doc", ".docx"}
EXCEL_EXTENSIONS  = {".xls", ".xlsx"}
PPT_EXTENSIONS    = {".ppt", ".pptx"}


# ─────────────────────────────────────────────
# Result object
# ─────────────────────────────────────────────

class ExtractionResult:
    def __init__(self):
        self.text: str = ""
        self.page_count: int = 0
        self.word_count: int = 0
        self.metadata: dict = {}
        self.title: Optional[str] = None
        self.author: Optional[str] = None
        self.subject: Optional[str] = None
        self.language: Optional[str] = None
        self.pages: list[dict] = []
        self.file_type: str = "unknown"
        self.error: Optional[str] = None
        self.success: bool = False

    def to_dict(self) -> dict:
        return {
            "text": self.text,
            "page_count": self.page_count,
            "word_count": self.word_count,
            "metadata": self.metadata,
            "title": self.title,
            "author": self.author,
            "subject": self.subject,
            "language": self.language,
            "pages": self.pages,
            "file_type": self.file_type,
            "success": self.success,
        }


# ─────────────────────────────────────────────
# Universal Extractor Service
# ─────────────────────────────────────────────

class DocumentExtractorService:

    def __init__(self):
        self.upload_dir = Path(settings.UPLOAD_DIR)
        self.upload_dir.mkdir(parents=True, exist_ok=True)

    # ── Save uploaded file ──────────────────────────────────────────────────

    async def save_uploaded_file(self, file_content: bytes, original_filename: str) -> dict:
        """Simpan file ke disk dengan nama unik. Return dict info file."""
        ext = Path(original_filename).suffix.lower()

        if ext not in SUPPORTED_EXTENSIONS:
            supported = ", ".join(sorted(SUPPORTED_EXTENSIONS.keys()))
            raise ValueError(
                f"Format '{ext}' tidak didukung. "
                f"Format yang didukung: {supported}"
            )

        file_size = len(file_content)
        if file_size > settings.max_file_size_bytes:
            raise ValueError(
                f"File terlalu besar ({file_size / 1024 / 1024:.1f} MB). "
                f"Maksimum {settings.MAX_FILE_SIZE_MB} MB."
            )

        unique_name = (
            f"{uuid.uuid4().hex}_{datetime.now(timezone.utc).strftime('%Y%m%d%H%M%S')}{ext}"
        )
        file_path = self.upload_dir / unique_name

        async with aiofiles.open(file_path, "wb") as f:
            await f.write(file_content)

        logger.info(f"Saved: {file_path} ({file_size} bytes)")

        return {
            "filename":          unique_name,
            "original_filename": original_filename,
            "file_path":         str(file_path),
            "file_size_bytes":   file_size,
            "mime_type":         SUPPORTED_EXTENSIONS[ext],
        }

    # ── Main extract dispatcher ─────────────────────────────────────────────

    def extract(self, file_path: str) -> ExtractionResult:
        """
        Entry point: deteksi format lalu dispatch ke handler yang sesuai.
        Jalankan di thread pool dari async context.
        """
        ext = Path(file_path).suffix.lower()
        result = ExtractionResult()
        result.file_type = ext.lstrip(".")

        try:
            if ext in PDF_EXTENSIONS:
                self._extract_pdf(file_path, result)
            elif ext in WORD_EXTENSIONS:
                self._extract_word(file_path, result)
            elif ext in EXCEL_EXTENSIONS:
                self._extract_excel(file_path, result)
            elif ext in PPT_EXTENSIONS:
                self._extract_pptx(file_path, result)
            elif ext in TEXT_EXTENSIONS:
                self._extract_text(file_path, result)
            elif ext in IMAGE_EXTENSIONS:
                self._extract_image(file_path, result)
            else:
                result.error = f"Format '{ext}' tidak didukung."
                return result

            if result.text:
                result.word_count = len(result.text.split())
                result.language   = self._detect_language(result.text[:1000])
                result.success    = True
            else:
                result.success = False
                if not result.error:
                    result.error = "Tidak ada teks yang bisa diekstrak dari file ini."

        except ImportError as e:
            missing_lib = str(e).split("'")[1] if "'" in str(e) else str(e)
            result.error   = f"Library '{missing_lib}' belum terinstall. Jalankan: pip install {missing_lib}"
            result.success = False
            logger.error(f"Missing library for {ext}: {e}")
        except Exception as e:
            result.error   = f"Gagal mengekstrak '{Path(file_path).name}': {str(e)}"
            result.success = False
            logger.error(f"Extraction error for {file_path}: {e}", exc_info=True)

        return result

    # ── PDF ────────────────────────────────────────────────────────────────

    def _extract_pdf(self, file_path: str, result: ExtractionResult):
        import fitz  # PyMuPDF

        doc = fitz.open(file_path)
        try:
            meta             = doc.metadata or {}
            result.metadata  = meta
            result.title     = meta.get("title") or None
            result.author    = meta.get("author") or None
            result.subject   = meta.get("subject") or None
            result.page_count = doc.page_count

            parts = []
            for i in range(doc.page_count):
                page      = doc[i]
                page_text = page.get_text("text")
                parts.append(page_text)
                result.pages.append({
                    "page_number": i + 1,
                    "char_count":  len(page_text),
                    "word_count":  len(page_text.split()),
                    "has_images":  len(page.get_images()) > 0,
                })

            result.text = "\n\n".join(parts).strip()
        finally:
            doc.close()

    # ── Word (.docx & .doc) ────────────────────────────────────────────────

    def _extract_word(self, file_path: str, result: ExtractionResult):
        ext = Path(file_path).suffix.lower()

        if ext == ".docx":
            from docx import Document  # python-docx
            doc = Document(file_path)

            # Core properties
            props         = doc.core_properties
            result.title  = props.title or None
            result.author = props.author or None

            parts = []
            for para in doc.paragraphs:
                if para.text.strip():
                    parts.append(para.text)

            # Tabel
            for table in doc.tables:
                rows = []
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if cells:
                        rows.append(" | ".join(cells))
                if rows:
                    parts.append("\n".join(rows))

            result.text       = "\n".join(parts)
            result.page_count = 1  # python-docx tidak expose page count

        else:  # .doc (legacy binary format)
            # Coba antiword atau textract sebagai fallback
            try:
                import subprocess
                proc = subprocess.run(
                    ["antiword", file_path],
                    capture_output=True, text=True, timeout=30
                )
                if proc.returncode == 0:
                    result.text = proc.stdout
                else:
                    raise RuntimeError("antiword gagal")
            except (FileNotFoundError, RuntimeError):
                # Fallback: baca raw bytes, ambil teks printable
                with open(file_path, "rb") as f:
                    raw = f.read()
                # Ambil string ASCII printable dari binary .doc
                text = raw.decode("latin-1", errors="ignore")
                printable = re.sub(r"[^\x20-\x7E\n\r\t]", " ", text)
                # Bersihkan spasi berlebih
                result.text = re.sub(r" {3,}", " ", printable).strip()
                logger.warning(f".doc fallback raw-text used for {file_path}")

            result.page_count = 1

    # ── Excel (.xlsx & .xls) ──────────────────────────────────────────────

    def _extract_excel(self, file_path: str, result: ExtractionResult):
        import openpyxl
        from openpyxl import load_workbook

        ext = Path(file_path).suffix.lower()

        if ext == ".xlsx":
            wb = load_workbook(file_path, read_only=True, data_only=True)
        else:
            # .xls — gunakan xlrd lalu convert representation ke teks
            try:
                import xlrd
                book   = xlrd.open_workbook(file_path)
                parts  = []
                for sheet_idx in range(book.nsheets):
                    sheet = book.sheet_by_index(sheet_idx)
                    parts.append(f"=== Sheet: {sheet.name} ===")
                    for row_idx in range(sheet.nrows):
                        row_vals = [str(sheet.cell_value(row_idx, col)).strip()
                                    for col in range(sheet.ncols)]
                        non_empty = [v for v in row_vals if v and v != "0.0" or v == "0"]
                        if non_empty:
                            parts.append(" | ".join(row_vals))
                result.text       = "\n".join(parts)
                result.page_count = book.nsheets
                return
            except ImportError:
                raise ImportError("xlrd")  # akan ditangkap di caller

        # xlsx path
        parts = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            parts.append(f"=== Sheet: {sheet_name} ===")
            for row in ws.iter_rows(values_only=True):
                cells = [str(c).strip() if c is not None else "" for c in row]
                non_empty = [c for c in cells if c]
                if non_empty:
                    parts.append(" | ".join(cells))

        result.text       = "\n".join(parts)
        result.page_count = len(wb.sheetnames)
        result.title      = Path(file_path).stem
        wb.close()

    # ── PowerPoint (.pptx & .ppt) ────────────────────────────────────────

    def _extract_pptx(self, file_path: str, result: ExtractionResult):
        ext = Path(file_path).suffix.lower()

        if ext == ".pptx":
            from pptx import Presentation
            from pptx.util import Pt

            prs   = Presentation(file_path)
            parts = []

            for slide_num, slide in enumerate(prs.slides, 1):
                slide_parts = [f"=== Slide {slide_num} ==="]
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_parts.append(shape.text.strip())
                    # Tabel di dalam slide
                    if shape.has_table:
                        for row in shape.table.rows:
                            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                            if cells:
                                slide_parts.append(" | ".join(cells))
                parts.append("\n".join(slide_parts))

            result.text       = "\n\n".join(parts)
            result.page_count = len(prs.slides)

            # Core properties
            props         = prs.core_properties
            result.title  = props.title or None
            result.author = props.author or None

        else:  # .ppt legacy
            result.error = (
                "Format .ppt (PowerPoint lama) tidak didukung secara langsung. "
                "Silakan convert ke .pptx terlebih dahulu."
            )

    # ── Plain text / CSV / TSV ────────────────────────────────────────────

    def _extract_text(self, file_path: str, result: ExtractionResult):
        ext = Path(file_path).suffix.lower()

        # Coba beberapa encoding
        encodings = ["utf-8", "utf-8-sig", "latin-1", "cp1252"]
        text      = None
        for enc in encodings:
            try:
                with open(file_path, "r", encoding=enc) as f:
                    text = f.read()
                break
            except UnicodeDecodeError:
                continue

        if text is None:
            with open(file_path, "rb") as f:
                raw  = f.read()
            text = raw.decode("latin-1", errors="replace")

        if ext == ".csv":
            # Format ulang CSV jadi teks yang mudah dibaca AI
            import csv, io
            reader  = csv.reader(io.StringIO(text))
            rows    = list(reader)
            if rows:
                header = rows[0]
                lines  = [" | ".join(header)]
                lines += [" | ".join(row) for row in rows[1:] if any(c.strip() for c in row)]
                text   = "\n".join(lines)
            result.page_count = 1

        elif ext == ".tsv":
            import csv, io
            reader  = csv.reader(io.StringIO(text), delimiter="\t")
            rows    = list(reader)
            if rows:
                lines = [" | ".join(row) for row in rows if any(c.strip() for c in row)]
                text  = "\n".join(lines)
            result.page_count = 1

        else:
            result.page_count = text.count("\n") // 50 + 1  # estimasi halaman

        result.text  = text.strip()
        result.title = Path(file_path).stem

    # ── Gambar / Foto / Struk ─────────────────────────────────────────────

    def _extract_image(self, file_path: str, result: ExtractionResult):
        """
        OCR gambar menggunakan pytesseract + PIL.
        Cocok untuk struk belanja, nota, foto dokumen, dsb.
        """
        try:
            import pytesseract
            from PIL import Image, ImageEnhance, ImageFilter
        except ImportError:
            raise ImportError("pytesseract")

        img = Image.open(file_path)

        # Pre-processing untuk meningkatkan akurasi OCR
        img = self._preprocess_image_for_ocr(img)

        # OCR dengan konfigurasi optimal untuk dokumen
        custom_config = r"--oem 3 --psm 6 -l ind+eng"
        try:
            text = pytesseract.image_to_string(img, config=custom_config)
        except pytesseract.TesseractNotFoundError:
            result.error = (
                "Tesseract OCR tidak ditemukan. "
                "Install Tesseract: https://github.com/tesseract-ocr/tesseract "
                "dan pastikan ada di PATH."
            )
            return

        result.text       = text.strip()
        result.page_count = 1
        result.title      = Path(file_path).stem

        # Metadata gambar
        try:
            result.metadata = {
                "width":  img.width,
                "height": img.height,
                "mode":   img.mode,
            }
        except Exception:
            pass

    def _preprocess_image_for_ocr(self, img):
        """
        Pre-processing gambar supaya OCR lebih akurat:
        - Konversi ke grayscale
        - Sharpening
        - Sedikit denoise
        Tidak aggressive (hindari merusak karakter)
        """
        from PIL import Image, ImageEnhance, ImageFilter

        # Konversi ke RGB dulu kalau ada alpha channel (RGBA/PNG)
        if img.mode in ("RGBA", "P"):
            img = img.convert("RGB")

        # Grayscale
        img = img.convert("L")

        # Sharpen
        img = img.filter(ImageFilter.SHARPEN)

        # Tingkatkan kontras sedikit
        enhancer = ImageEnhance.Contrast(img)
        img       = enhancer.enhance(1.5)

        return img

    # ── Language detection ────────────────────────────────────────────────

    def _detect_language(self, text: str) -> str:
        text_lower = text.lower()
        id_words   = ["yang", "dan", "di", "ini", "itu", "untuk", "dengan", "tidak", "dari", "ke"]
        en_words   = ["the", "and", "is", "are", "was", "were", "for", "with", "this", "that"]
        id_count   = sum(1 for w in id_words if f" {w} " in text_lower)
        en_count   = sum(1 for w in en_words if f" {w} " in text_lower)
        if id_count > en_count:
            return "id"
        elif en_count > 0:
            return "en"
        return "unknown"

    # ── Utilities ─────────────────────────────────────────────────────────

    def get_text_preview(self, text: str, max_chars: int = 500) -> str:
        if not text:
            return ""
        preview = text[:max_chars].strip()
        if len(text) > max_chars:
            preview += "..."
        return preview

    def chunk_text(self, text: str, chunk_size: int = 3000, overlap: int = 200) -> list[str]:
        if len(text) <= chunk_size:
            return [text]
        chunks, start = [], 0
        while start < len(text):
            end   = start + chunk_size
            chunk = text[start:end]
            if end < len(text):
                last_period = chunk.rfind(". ")
                if last_period > chunk_size // 2:
                    chunk = chunk[:last_period + 1]
                    end   = start + last_period + 1
            chunks.append(chunk.strip())
            start = end - overlap
        return chunks

    async def delete_file(self, file_path: str) -> bool:
        try:
            path = Path(file_path)
            if path.exists():
                path.unlink()
                logger.info(f"Deleted file: {file_path}")
            return True
        except Exception as e:
            logger.error(f"Error deleting file {file_path}: {e}")
            return False


# Singleton — gantikan pdf_extractor di seluruh codebase
document_extractor = DocumentExtractorService()

# Alias backward-compat supaya import lama (pdf_extractor) tidak langsung error
pdf_extractor = document_extractor